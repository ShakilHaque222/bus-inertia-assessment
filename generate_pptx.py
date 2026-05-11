"""Generate comprehensive PowerPoint presentation for bus-level inertia assessment."""

import os, sys
sys.path.insert(0, os.path.dirname(__file__))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as In
import io, math

# ── Paths ──────────────────────────────────────────────────────────────────────
RESULTS  = os.path.join(os.path.dirname(__file__), 'results')
PPTX_OUT = os.path.join(RESULTS, 'Bus_Inertia_Presentation.pptx')

# ── Colour palette ─────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0d, 0x1b, 0x4b)
BLUE   = RGBColor(0x1a, 0x3a, 0x8c)
LTBLUE = RGBColor(0x3a, 0x6b, 0xbf)
ACCENT = RGBColor(0xe8, 0x40, 0x40)
GOLD   = RGBColor(0xc8, 0xa0, 0x00)
LGREY  = RGBColor(0xe8, 0xe8, 0xf0)
WHITE  = RGBColor(0xff, 0xff, 0xff)
DGREY  = RGBColor(0x4a, 0x4a, 0x5a)
CYAN   = RGBColor(0x00, 0xcc, 0xcc)
GREEN  = RGBColor(0x44, 0xcc, 0x88)
ORANGE = RGBColor(0xff, 0x8c, 0x00)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helper functions ───────────────────────────────────────────────────────────
def rgb(r, g, b): return RGBColor(r, g, b)

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.line.fill.background()  # no line by default
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
        else:
            shape.line.fill.background()
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_name='Calibri', font_size=14, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, wrap=True, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_text(slide, lines, left, top, width, height,
                    font_name='Calibri', font_size=13, color=WHITE,
                    bullet_color=None, line_spacing=1.2):
    from pptx.oxml.ns import qn
    from pptx.util import Pt
    from lxml import etree

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (bullet, line) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(3)
        if bullet:
            run0 = p.add_run()
            run0.text = '▸  '
            run0.font.name = font_name
            run0.font.size = Pt(font_size)
            run0.font.color.rgb = bullet_color or GOLD
            run0.font.bold = True
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox


def add_image(slide, fname, left, top, width=None, height=None):
    path = os.path.join(RESULTS, fname)
    if not os.path.exists(path):
        return None
    if width and height:
        return slide.shapes.add_picture(path, left, top, width, height)
    elif width:
        return slide.shapes.add_picture(path, left, top, width=width)
    elif height:
        return slide.shapes.add_picture(path, left, top, height=height)
    else:
        return slide.shapes.add_picture(path, left, top)


def header_bar(slide, title, subtitle=None):
    """Dark navy header with accent stripe."""
    # Main bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.15), fill_color=NAVY)
    # Accent stripe
    add_rect(slide, 0, Inches(1.15), SLIDE_W, Inches(0.04), fill_color=ACCENT)
    # Left accent bar
    add_rect(slide, 0, 0, Inches(0.08), SLIDE_H, fill_color=LTBLUE)

    add_text(slide, title,
             Inches(0.2), Inches(0.12), Inches(10), Inches(0.65),
             font_name='Calibri', font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.2), Inches(0.72), Inches(11), Inches(0.38),
                 font_name='Calibri', font_size=14, italic=True,
                 color=RGBColor(0xaa, 0xbb, 0xdd))


def slide_number(slide, n, total):
    add_text(slide, f'{n} / {total}',
             Inches(12.3), Inches(7.15), Inches(0.9), Inches(0.3),
             font_size=9, color=DGREY, align=PP_ALIGN.RIGHT)


def footer_bar(slide, text='Ghosh, Isbeih & El Moursi | IEEE Trans. Power Del. 2023'):
    add_rect(slide, 0, Inches(7.2), SLIDE_W, Inches(0.3), fill_color=LGREY)
    add_text(slide, text, Inches(0.15), Inches(7.22), Inches(11), Inches(0.25),
             font_size=8, color=DGREY)


def eq_box(slide, equation, left, top, width, height, font_size=14):
    add_rect(slide, left, top, width, height,
             fill_color=RGBColor(0xee, 0xf2, 0xff),
             line_color=LTBLUE, line_width=Pt(1))
    add_text(slide, equation, left + Inches(0.1), top + Inches(0.06),
             width - Inches(0.2), height - Inches(0.1),
             font_name='Consolas', font_size=font_size, bold=True,
             color=NAVY, align=PP_ALIGN.CENTER)


def badge(slide, text, left, top, width, height, bg=NAVY, fg=WHITE, font_size=11):
    add_rect(slide, left, top, width, height, fill_color=bg)
    add_text(slide, text, left, top, width, height,
             font_size=font_size, bold=True, color=fg, align=PP_ALIGN.CENTER)


# ── Slide builders ─────────────────────────────────────────────────────────────

def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Full background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=NAVY)
    # Accent top
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill_color=ACCENT)
    # Accent bottom
    add_rect(slide, 0, Inches(7.38), SLIDE_W, Inches(0.12), fill_color=LTBLUE)
    # Left bar
    add_rect(slide, 0, 0, Inches(0.18), SLIDE_H, fill_color=LTBLUE)
    # Right bar
    add_rect(slide, Inches(13.15), 0, Inches(0.18), SLIDE_H, fill_color=LTBLUE)

    # Central content panel
    add_rect(slide, Inches(0.5), Inches(0.8), Inches(12.33), Inches(5.9),
             fill_color=RGBColor(0x10, 0x22, 0x5a))

    # Title
    add_text(slide, 'Bus-Level Inertia Assessment',
             Inches(0.6), Inches(1.0), Inches(12.1), Inches(0.85),
             font_name='Calibri', font_size=36, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(slide, 'for Converter-Dominated Power Systems',
             Inches(0.6), Inches(1.75), Inches(12.1), Inches(0.65),
             font_name='Calibri', font_size=28, bold=False, color=LGREY,
             align=PP_ALIGN.CENTER)

    # Divider line
    add_rect(slide, Inches(1.5), Inches(2.52), Inches(10.33), Inches(0.02),
             fill_color=GOLD)

    # Reference
    add_text(slide, 'Based on:',
             Inches(0.6), Inches(2.65), Inches(12.1), Inches(0.3),
             font_size=11, italic=True, color=RGBColor(0xaa, 0xbb, 0xdd),
             align=PP_ALIGN.CENTER)
    add_text(slide,
             'Ghosh, Isbeih & El Moursi — IEEE Trans. Power Delivery, Vol. 38, No. 4, Aug. 2023',
             Inches(0.6), Inches(2.95), Inches(12.1), Inches(0.4),
             font_name='Calibri', font_size=14, bold=True, color=GOLD,
             align=PP_ALIGN.CENTER)

    # Four systems badges
    systems = ['IEEE 14-Bus\n5 Generators', 'IEEE 39-Bus\n10 Generators',
               'IEEE 68-Bus\n16 Generators', 'IEEE 118-Bus\n54 Generators']
    colors_b = [LTBLUE, BLUE, ACCENT, RGBColor(0x22, 0x88, 0x44)]
    for k, (s, c) in enumerate(zip(systems, colors_b)):
        badge(slide, s, Inches(0.7 + k*3.08), Inches(3.55),
              Inches(2.75), Inches(0.75), bg=c, font_size=12)

    # Key equation
    add_rect(slide, Inches(1.0), Inches(4.5), Inches(11.33), Inches(0.65),
             fill_color=RGBColor(0x18, 0x30, 0x70))
    add_text(slide, 'H_B(t) = W_c(t) × H_G          Sys_str(t) = Σ H_B(t)',
             Inches(1.1), Inches(4.55), Inches(11.1), Inches(0.55),
             font_name='Consolas', font_size=16, bold=True, color=CYAN,
             align=PP_ALIGN.CENTER)

    # Author
    add_text(slide, 'Shakil Haque  |  shakil.haqueee@gmail.com  |  May 2026',
             Inches(0.6), Inches(5.35), Inches(12.1), Inches(0.35),
             font_size=11, color=RGBColor(0x88, 0x99, 0xbb), align=PP_ALIGN.CENTER)

    add_text(slide, 'Python Simulation  ·  Y-bus Algorithm  ·  IEEE 14 / 39 / 68 / 118-Bus',
             Inches(0.6), Inches(5.65), Inches(12.1), Inches(0.3),
             font_size=10, italic=True, color=RGBColor(0x66, 0x77, 0x99),
             align=PP_ALIGN.CENTER)


def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, 'Presentation Overview', 'Research structure and contents')
    footer_bar(slide)

    sections = [
        ('01', 'Background', 'Power system inertia, the energy transition, swing equation'),
        ('02', 'Existing Methods', 'L1–L6 approaches and their limitations (gap analysis)'),
        ('03', 'Research Novelty', 'N1–N6: augmented Y-bus, W_c matrix, H_B formula'),
        ('04', 'Methodology', 'Algorithm: Y-bus construction → partitioning → H_B computation'),
        ('05', 'Network Topology', 'IEEE 14/39/68/118-bus system structures'),
        ('06', 'Results', 'H_B distribution, frequency response, RES penetration'),
        ('07', 'Analysis', 'Physical interpretation, VI placement, critical buses'),
        ('08', 'Future Work', 'PMU integration, optimal VI control, GNN acceleration'),
        ('09', 'Conclusion', 'Key findings and contributions summary'),
    ]

    cols = 3
    w, h = Inches(4.15), Inches(0.82)
    for k, (num, title, desc) in enumerate(sections):
        col = k % cols
        row = k // cols
        lf = Inches(0.22) + col * Inches(4.37)
        tp = Inches(1.35) + row * Inches(0.95)

        add_rect(slide, lf, tp, w, h, fill_color=BLUE)
        add_rect(slide, lf, tp, Inches(0.52), h, fill_color=LTBLUE)
        add_text(slide, num, lf + Inches(0.04), tp + Inches(0.2),
                 Inches(0.48), Inches(0.4), font_size=16, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, lf + Inches(0.58), tp + Inches(0.06),
                 Inches(3.5), Inches(0.35), font_size=12, bold=True, color=WHITE)
        add_text(slide, desc, lf + Inches(0.58), tp + Inches(0.42),
                 Inches(3.5), Inches(0.35), font_size=9, italic=True,
                 color=LGREY)


def slide_background(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, '1. Background & Motivation',
                'Why inertia matters in converter-dominated power systems')
    footer_bar(slide)

    # Left column: problem
    add_rect(slide, Inches(0.22), Inches(1.3), Inches(6.2), Inches(5.9),
             fill_color=RGBColor(0x12, 0x1e, 0x48))

    add_text(slide, 'The Inertia Problem', Inches(0.3), Inches(1.35),
             Inches(6.0), Inches(0.42), font_size=15, bold=True, color=GOLD)

    problems = [
        (True, 'Synchronous generators store kinetic energy (H [s])'),
        (True, 'RES (wind/solar) have NO inherent inertia'),
        (True, 'As RES replaces generators: system inertia falls'),
        (True, 'Result: faster frequency drop after disturbances'),
        (True, 'Higher ROCOF → risk of UFLS and blackouts'),
        (False, ''),
        (True, 'Traditional COI model: ONE H_sys for whole grid'),
        (True, 'But frequency varies spatially across the network!'),
        (True, 'Different buses experience different ROCOF values'),
        (True, 'Need bus-level inertia for targeted VI placement'),
    ]
    add_bullet_text(slide, problems, Inches(0.3), Inches(1.82),
                    Inches(6.0), Inches(4.8), font_size=11.5, color=WHITE)

    # Right column: swing equation + visual
    add_rect(slide, Inches(6.6), Inches(1.3), Inches(6.5), Inches(2.65),
             fill_color=RGBColor(0x12, 0x1e, 0x48))
    add_text(slide, 'Swing Equation', Inches(6.7), Inches(1.35),
             Inches(6.3), Inches(0.4), font_size=14, bold=True, color=GOLD)
    eq_box(slide, '2H/ω₀ · dω/dt = Pm − Pe − D·Δω',
           Inches(6.7), Inches(1.82), Inches(6.2), Inches(0.52), font_size=13)
    eq_box(slide, 'ROCOF = −f₀·ΔP / (2·H_sys)  [Hz/s]',
           Inches(6.7), Inches(2.42), Inches(6.2), Inches(0.52), font_size=13)

    add_text(slide, '↑ H_sys  →  ↓ ROCOF  →  ↑ Frequency Stability',
             Inches(6.7), Inches(3.02), Inches(6.3), Inches(0.35),
             font_size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # Key insight box
    add_rect(slide, Inches(6.6), Inches(4.1), Inches(6.5), Inches(3.1),
             fill_color=RGBColor(0x12, 0x1e, 0x48))
    add_text(slide, 'Key Insight', Inches(6.7), Inches(4.15),
             Inches(6.3), Inches(0.38), font_size=14, bold=True, color=GOLD)
    points = [
        (True, 'High H_B bus  →  slow frequency drop, deeper nadir'),
        (True, 'Low H_B bus   →  fast frequency drop, near UFLS'),
        (True, 'VI placement should target lowest-H_B buses'),
        (True, 'Bus-specific relay settings prevent false trips'),
        (True, 'RES forecast + H_B enables predictive control'),
    ]
    add_bullet_text(slide, points, Inches(6.7), Inches(4.58),
                    Inches(6.2), Inches(2.5), font_size=11.5, color=WHITE)


def slide_existing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, '2. Existing Approaches (L1–L6)',
                'Current methods and their critical limitations')
    footer_bar(slide)

    methods = [
        ('L1', 'PMU ROCOF',   'System-wide only\nNeeds disturbance\nHigh noise sensitivity'),
        ('L2', 'Ambient Data','High variance\nLong windows (min)\nNo spatial resolution'),
        ('L3', 'ML/Data-Driven','Black-box\nNeeds labelled data\nNo generalization'),
        ('L4', 'COI Model',  'All buses same H\nPhysically wrong\nNo VI guidance'),
        ('L5', 'Modal Analysis','Needs full model\nNo H_B per bus\nNot real-time'),
        ('L6', 'Zone-Based', 'Arbitrary zones\nIntra-zone blind\nManual update'),
    ]

    colors_m = [LTBLUE, BLUE, RGBColor(0x33, 0x66, 0xaa),
                RGBColor(0x55, 0x44, 0x99), BLUE, LTBLUE]

    for k, ((code, name, lims), mc) in enumerate(zip(methods, colors_m)):
        col = k % 3
        row = k // 3
        lf = Inches(0.22) + col * Inches(4.37)
        tp = Inches(1.32) + row * Inches(2.8)
        w, h = Inches(4.15), Inches(2.65)

        add_rect(slide, lf, tp, w, h, fill_color=mc)
        badge(slide, code, lf, tp, Inches(0.6), Inches(0.42),
              bg=NAVY, font_size=13)
        add_text(slide, name, lf + Inches(0.66), tp + Inches(0.06),
                 Inches(3.4), Inches(0.35), font_size=13, bold=True, color=WHITE)

        add_rect(slide, lf + Inches(0.08), tp + Inches(0.5),
                 Inches(3.98), Inches(0.03), fill_color=GOLD)

        for j, lim in enumerate(lims.split('\n')):
            add_text(slide, f'✗  {lim}',
                     lf + Inches(0.12), tp + Inches(0.62 + j * 0.57),
                     Inches(3.9), Inches(0.52),
                     font_size=11, color=LGREY)

    # Bottom note
    add_rect(slide, Inches(0.22), Inches(6.88), Inches(12.88), Inches(0.32),
             fill_color=RGBColor(0x1a, 0x3a, 0x8c))
    add_text(slide,
             'None of the above methods provide BUS-LEVEL inertia with physical interpretability and real-time capability.',
             Inches(0.32), Inches(6.9), Inches(12.68), Inches(0.28),
             font_size=11, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def slide_novelty(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, '3. Research Novelty (N1–N6)',
                'Six contributions of Ghosh, Isbeih & El Moursi (2023)')
    footer_bar(slide)

    novelties = [
        ('N1', 'Augmented Y-bus Partitioning',
         'Y partitioned into Y_gg, Y_gb,\nY_bg, Y_bb submatrices'),
        ('N2', 'Weighting Matrix W_c',
         'Electrical coupling from each\nload bus to each generator'),
        ('N3', 'Bus-Level Inertia H_B',
         'H_B(i) = Σⱼ W_c,ij · H_Gʲ\nUnique H per bus'),
        ('N4', 'System Strength Metric',
         'Sys_str = Σ H_B\nReal-time inertia tracking'),
        ('N5', 'Fault-Mode Adaptation',
         'v_b → 1 pu during fault\nContinuous operation'),
        ('N6', 'VI Placement Guidance',
         'argmin{H_B} identifies\nhighest-priority VI buses'),
    ]

    ncol = [ACCENT, GOLD, GREEN, CYAN, ORANGE, LTBLUE]
    for k, ((code, title, desc), nc) in enumerate(zip(novelties, ncol)):
        col = k % 3
        row = k // 3
        lf = Inches(0.22) + col * Inches(4.37)
        tp = Inches(1.32) + row * Inches(2.35)
        w, h = Inches(4.15), Inches(2.2)

        add_rect(slide, lf, tp, w, h, fill_color=NAVY)
        add_rect(slide, lf, tp, w, Inches(0.06), fill_color=nc)
        badge(slide, code, lf + Inches(0.08), tp + Inches(0.12),
              Inches(0.6), Inches(0.38), bg=nc, fg=NAVY, font_size=14)
        add_text(slide, title, lf + Inches(0.74), tp + Inches(0.15),
                 Inches(3.3), Inches(0.38), font_size=12, bold=True, color=WHITE)
        add_text(slide, desc, lf + Inches(0.12), tp + Inches(0.62),
                 Inches(3.9), Inches(1.45), font_size=11.5, color=LGREY)

    # Core equation box
    add_rect(slide, Inches(0.22), Inches(5.88), Inches(12.88), Inches(1.22),
             fill_color=RGBColor(0x10, 0x22, 0x5a))
    add_text(slide, 'Core Algorithm Equations',
             Inches(0.32), Inches(5.92), Inches(12.68), Inches(0.35),
             font_size=12, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    eq_box(slide,
           'W_c(t) = { −Y_bb⁻¹(t) × Y_bg(t) ÷ v_b(t) } ⊙ e_g(t)     →     H_B(t) = W_c(t) × H_G',
           Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.65), font_size=14)


def slide_methodology(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, '4. Algorithm Methodology',
                'Step-by-step computation of bus-level inertia')
    footer_bar(slide)

    steps = [
        ('①', 'Branch Data Input',
         'R, X, B/2 for all lines\nGenerator buses + H_G values'),
        ('②', 'Build Y-bus',
         'Y[i,j] = −yᵢⱼ (off-diag)\nY[i,i] = Σ yᵢⱼ + y_shunt'),
        ('③', 'Partition Y-bus',
         'Split into: Y_gg, Y_gb\nY_bg, Y_bb submatrices'),
        ('④', 'Invert Y_bb',
         'Y_bb⁻¹ ∈ ℂⁿᵇˣⁿᵇ\nO(nb³) complexity'),
        ('⑤', 'Compute W_c',
         'M = −Y_bb⁻¹ × Y_bg\nNorm + voltage scaling'),
        ('⑥', 'Bus Inertia H_B',
         'H_B = W_c × H_G\nOne value per bus [s]'),
    ]

    arrow_color = GOLD
    for k, (num, title, desc) in enumerate(steps):
        lf = Inches(0.22 + k * 2.2)
        tp = Inches(1.32)
        w, h = Inches(2.0), Inches(3.5)

        add_rect(slide, lf, tp, w, h, fill_color=BLUE)
        add_rect(slide, lf, tp, w, Inches(0.06), fill_color=LTBLUE)
        add_text(slide, num, lf, tp + Inches(0.12), w, Inches(0.55),
                 font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, lf + Inches(0.08), tp + Inches(0.72),
                 Inches(1.84), Inches(0.45), font_size=11, bold=True, color=GOLD)
        add_text(slide, desc, lf + Inches(0.08), tp + Inches(1.22),
                 Inches(1.84), Inches(2.15), font_size=10.5, color=LGREY)

        if k < 5:
            add_text(slide, '→', Inches(0.22 + (k+1)*2.2 - 0.22),
                     tp + Inches(1.3), Inches(0.25), Inches(0.45),
                     font_size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Y-bus partition diagram
    add_rect(slide, Inches(0.22), Inches(5.0), Inches(6.0), Inches(2.1),
             fill_color=RGBColor(0x12, 0x1e, 0x48))
    add_text(slide, 'Y-bus Partition (Augmented Admittance Matrix)',
             Inches(0.32), Inches(5.05), Inches(5.8), Inches(0.38),
             font_size=12, bold=True, color=GOLD)
    add_text(slide,
             '  [  Y_gg    Y_gb  ]   [  v_g  ]   [  i_g  ]\n'
             '  [              ] × [       ] = [       ]\n'
             '  [  Y_bg    Y_bb  ]   [  v_b  ]   [  i_b  ]',
             Inches(0.32), Inches(5.48), Inches(5.8), Inches(1.4),
             font_name='Consolas', font_size=11.5, color=CYAN)

    # Weighting formula
    add_rect(slide, Inches(6.4), Inches(5.0), Inches(6.7), Inches(2.1),
             fill_color=RGBColor(0x12, 0x1e, 0x48))
    add_text(slide, 'Weighting Matrix — Physical Meaning',
             Inches(6.5), Inches(5.05), Inches(6.5), Inches(0.38),
             font_size=12, bold=True, color=GOLD)
    points = [
        (True, 'High |W_c(i,j)| → bus i electrically close to gen j'),
        (True, 'Low  |W_c(i,j)| → bus i far from gen j'),
        (True, 'e_g(j) = 0 when gen j is tripped (RES replaces it)'),
        (True, 'H_B decreases as generators are displaced by RES'),
    ]
    add_bullet_text(slide, points, Inches(6.5), Inches(5.48),
                    Inches(6.5), Inches(1.5), font_size=10.5, color=WHITE)


def slide_topology(prs, system_name, n_bus, n_gen, n_branch, fig_file,
                   gen_buses_str, key_result, slide_num=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, f'5. Network Topology — {system_name}',
                f'{n_bus} buses  |  {n_gen} generators  |  {n_branch} branches')
    footer_bar(slide)

    # Network image (left 60%)
    add_image(slide, fig_file, Inches(0.22), Inches(1.3), width=Inches(8.0))

    # Right panel
    add_rect(slide, Inches(8.4), Inches(1.3), Inches(4.7), Inches(5.9),
             fill_color=RGBColor(0x12, 0x1e, 0x48))
    add_text(slide, 'System Parameters',
             Inches(8.5), Inches(1.35), Inches(4.5), Inches(0.38),
             font_size=13, bold=True, color=GOLD)

    params = [
        f'Buses:          {n_bus}',
        f'Generators:     {n_gen}',
        f'Branches:       {n_branch}',
        f'Base MVA:       100 MVA',
        f'Gen buses: {gen_buses_str}',
    ]
    for j, p in enumerate(params):
        add_text(slide, p, Inches(8.5), Inches(1.82 + j * 0.44),
                 Inches(4.5), Inches(0.38), font_name='Consolas',
                 font_size=10.5, color=WHITE)

    add_rect(slide, Inches(8.5), Inches(4.1), Inches(4.4), Inches(0.03),
             fill_color=GOLD)
    add_text(slide, 'Key Result',
             Inches(8.5), Inches(4.18), Inches(4.5), Inches(0.35),
             font_size=12, bold=True, color=GOLD)
    add_text(slide, key_result, Inches(8.5), Inches(4.58),
             Inches(4.5), Inches(2.6), font_size=11, color=LGREY)

    # Legend
    add_rect(slide, Inches(0.22), Inches(6.88), Inches(4.0), Inches(0.32),
             fill_color=NAVY)
    add_text(slide, '■  Red squares = Generator buses     ●  Blue circles = Load buses',
             Inches(0.32), Inches(6.9), Inches(3.8), Inches(0.28),
             font_size=9, color=LGREY)


def slide_results_bar(prs, system_name, bar_fig, freq_fig):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, f'6. Results — {system_name}',
                'Bus-level inertia distribution and frequency response')
    footer_bar(slide)

    # Two images side by side
    add_image(slide, bar_fig,  Inches(0.22), Inches(1.32), width=Inches(6.4))
    add_image(slide, freq_fig, Inches(6.8),  Inches(1.32), width=Inches(6.3))

    add_text(slide, 'Bus Inertia Distribution (H_B)',
             Inches(0.22), Inches(6.85), Inches(6.4), Inches(0.28),
             font_size=9.5, italic=True, color=DGREY, align=PP_ALIGN.CENTER)
    add_text(slide, 'Frequency Response (ΔP = 0.1 pu)',
             Inches(6.8), Inches(6.85), Inches(6.3), Inches(0.28),
             font_size=9.5, italic=True, color=DGREY, align=PP_ALIGN.CENTER)


def slide_res_comparison(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, '6. Results — RES Penetration & System Comparison',
                'Impact of renewable energy integration on system inertia')
    footer_bar(slide)

    add_image(slide, 'res_impact_ieee39bus.png', Inches(0.22), Inches(1.32), width=Inches(6.4))
    add_image(slide, 'system_comparison_all.png', Inches(6.8), Inches(1.32), width=Inches(6.3))

    add_text(slide, 'IEEE 39-Bus: System Strength vs. RES Penetration',
             Inches(0.22), Inches(6.85), Inches(6.4), Inches(0.28),
             font_size=9.5, italic=True, color=DGREY, align=PP_ALIGN.CENTER)
    add_text(slide, 'Cross-System Inertia Comparison (All IEEE Test Systems)',
             Inches(6.8), Inches(6.85), Inches(6.3), Inches(0.28),
             font_size=9.5, italic=True, color=DGREY, align=PP_ALIGN.CENTER)


def slide_wc_heatmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, '6. Results — Weighting Matrix W_c',
                'Electrical coupling between load buses and generators')
    footer_bar(slide)

    add_image(slide, 'weighting_matrix_ieee39bus.png', Inches(0.22), Inches(1.32), width=Inches(7.5))

    # Interpretation panel
    add_rect(slide, Inches(7.9), Inches(1.32), Inches(5.2), Inches(5.9),
             fill_color=RGBColor(0x12, 0x1e, 0x48))
    add_text(slide, 'How to Read This Heatmap',
             Inches(8.0), Inches(1.38), Inches(5.0), Inches(0.38),
             font_size=13, bold=True, color=GOLD)

    interp = [
        (True, 'Each row = one load bus'),
        (True, 'Each column = one generator bus'),
        (True, 'Colour intensity = coupling strength W_c'),
        (True, 'Warm (yellow/red) = strong coupling'),
        (True, 'Cool (dark) = weak coupling'),
        (False, ''),
        (True, 'Bus 39 (H=500s) column is brightest'),
        (True, '→ dominates inertia for buses 1-15'),
        (False, ''),
        (True, 'Bus 20 has low coupling to all gens'),
        (True, '→ lowest H_B, highest ROCOF risk'),
        (False, ''),
        (True, 'Identify which gen "feeds" which bus'),
        (True, '→ guides generator retirement planning'),
    ]
    add_bullet_text(slide, interp, Inches(8.0), Inches(1.82),
                    Inches(5.0), Inches(5.0), font_size=11, color=WHITE)


def slide_analysis(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, '7. Analysis — Key Physical Insights',
                'What the results tell us about grid inertia')
    footer_bar(slide)

    findings = [
        ('Spatial Variation is Extreme',
         ACCENT,
         ['IEEE 39-Bus: H_B ranges 26 s → 374 s (14× ratio)',
          'Uniform COI model completely misrepresents this',
          'Network topology determines inertia distribution']),
        ('ROCOF Matches H_B Precisely',
         GOLD,
         ['Bus with 14× less H_B has 14× higher ROCOF',
          'Validated by swing equation simulation',
          'Enables bus-specific relay threshold settings']),
        ('RES Integration Critical Threshold',
         CYAN,
         ['IEEE 14-Bus: >60% RES → >60% strength loss',
          'Non-linear: last large-H gen removal worst case',
          'Early warning from H_B monitoring essential']),
        ('Generator 39 Dominates (39-Bus)',
         GREEN,
         ['H=500s "infinite bus" determines inertia for 40% of buses',
          'Its retirement would devastate Bus 1–15 inertia',
          'Identified by W_c heatmap column analysis']),
    ]

    for k, (title, tc, points) in enumerate(findings):
        col = k % 2
        row = k // 2
        lf = Inches(0.22) + col * Inches(6.55)
        tp = Inches(1.32) + row * Inches(2.95)
        w, h = Inches(6.35), Inches(2.78)

        add_rect(slide, lf, tp, w, h, fill_color=NAVY)
        add_rect(slide, lf, tp, w, Inches(0.06), fill_color=tc)
        add_text(slide, title, lf + Inches(0.12), tp + Inches(0.12),
                 Inches(6.1), Inches(0.4), font_size=13, bold=True, color=tc)
        pts = [(True, p) for p in points]
        add_bullet_text(slide, pts, lf + Inches(0.12), tp + Inches(0.6),
                        Inches(6.1), Inches(2.1), font_size=11.5,
                        color=WHITE, bullet_color=tc)


def slide_future(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, '8. Future Work',
                'Extending the bus-level inertia framework')
    footer_bar(slide)

    items = [
        ('F1', 'PMU Integration',    'Real-time H_B at PMU\nsampling rate (30–120 Hz)', LTBLUE),
        ('F2', 'Optimal VI Control', 'Closed-loop VI from\ngrid-forming inverters', BLUE),
        ('F3', 'Probabilistic H_B',  'E[H_B], Var[H_B] under\nRES output uncertainty', ACCENT),
        ('F4', '3-Phase Systems',    'Unbalanced distribution\nnetwork extension', GREEN),
        ('F5', 'Converter Inertia',  'VSM-based programmable\nH_G(t) optimisation', CYAN),
        ('F6', 'HIL Validation',     'RTDS / OPAL-RT testing\nwith physical relays', ORANGE),
        ('F7', 'HVDC Extension',     'Multi-area with DC ties\nand converter control', LTBLUE),
        ('F8', 'GNN Acceleration',   'Graph Neural Network to\npredict H_B at >1000 buses', GOLD),
    ]

    for k, (code, title, desc, ic) in enumerate(items):
        col = k % 4
        row = k // 4
        lf = Inches(0.22) + col * Inches(3.27)
        tp = Inches(1.32) + row * Inches(2.5)
        w, h = Inches(3.1), Inches(2.35)

        add_rect(slide, lf, tp, w, h, fill_color=NAVY)
        add_rect(slide, lf, tp, Inches(0.52), h, fill_color=ic)
        add_text(slide, code, lf + Inches(0.06), tp + Inches(0.85),
                 Inches(0.42), Inches(0.48), font_size=14, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, title, lf + Inches(0.6), tp + Inches(0.12),
                 Inches(2.42), Inches(0.45), font_size=12, bold=True, color=ic)
        add_text(slide, desc, lf + Inches(0.6), tp + Inches(0.62),
                 Inches(2.42), Inches(1.6), font_size=11, color=LGREY)

    # Roadmap
    add_rect(slide, Inches(0.22), Inches(6.8), Inches(12.88), Inches(0.4),
             fill_color=RGBColor(0x10, 0x22, 0x5a))
    add_text(slide,
             'Roadmap:  Phase 1 (0–12 mo): PMU + HIL  →  Phase 2 (12–24 mo): VI Control + Probabilistic  →  Phase 3 (24–36 mo): HVDC + GNN',
             Inches(0.32), Inches(6.83), Inches(12.68), Inches(0.35),
             font_size=10.5, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def slide_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header_bar(slide, '9. Conclusion',
                'Summary of contributions and key findings')
    footer_bar(slide)

    conclusions = [
        (ACCENT, 'Spatial Inertia',
         'H_B varies up to 14× across a single network — uniform COI model is insufficient'),
        (GOLD, 'Y-bus Algorithm',
         'Y_bb⁻¹ × Y_bg partitioning provides physical, interpretable, real-time H_B'),
        (GREEN, 'ROCOF Validation',
         'Swing equation confirms H_B → ROCOF correlation validated across all 4 systems'),
        (CYAN, 'RES Impact Tracked',
         'System strength decreases monotonically — framework enables early warning'),
        (ORANGE, 'VI Placement Guided',
         'Lowest H_B buses = highest-priority sites for grid-forming inverter deployment'),
        (LTBLUE, 'Computationally Efficient',
         'Single matrix inversion O(nb³) — scalable to production EMS integration'),
    ]

    for k, (tc, title, text) in enumerate(conclusions):
        tp = Inches(1.35) + k * Inches(0.87)
        add_rect(slide, Inches(0.22), tp, Inches(12.88), Inches(0.8), fill_color=NAVY)
        add_rect(slide, Inches(0.22), tp, Inches(0.06), Inches(0.8), fill_color=tc)
        add_text(slide, title, Inches(0.36), tp + Inches(0.08),
                 Inches(3.2), Inches(0.35), font_size=12, bold=True, color=tc)
        add_text(slide, text, Inches(3.65), tp + Inches(0.15),
                 Inches(9.3), Inches(0.5), font_size=11.5, color=WHITE)

    # Final quote
    add_rect(slide, Inches(0.22), Inches(6.68), Inches(12.88), Inches(0.52),
             fill_color=RGBColor(0x10, 0x22, 0x5a))
    add_text(slide,
             '"Bus-level inertia assessment is the foundation for proactive frequency stability management in the converter-dominated grid of the future."',
             Inches(0.32), Inches(6.72), Inches(12.68), Inches(0.42),
             font_size=11, italic=True, color=GOLD, align=PP_ALIGN.CENTER)


# ── main ───────────────────────────────────────────────────────────────────────
def main():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    print('Building PowerPoint slides...')
    slide_cover(prs);       print('  Slide 1: Cover')
    slide_overview(prs);    print('  Slide 2: Overview')
    slide_background(prs);  print('  Slide 3: Background')
    slide_existing(prs);    print('  Slide 4: Existing Methods')
    slide_novelty(prs);     print('  Slide 5: Novelty N1-N6')
    slide_methodology(prs); print('  Slide 6: Methodology')

    # 4 topology slides
    topo_data = [
        ('IEEE 14-Bus', 14, 5, 20, 'topology_ieee14bus.png',
         '1, 2, 3, 6, 8',
         'Max H_B: Bus 1 (10.3 s)\nMin H_B: Bus 8 (4.0 s)\nSys strength: 75.3 s\nAvg H_B: 5.4 s\nH_B ratio (max/min): 2.6×'),
        ('IEEE 39-Bus', 39, 10, 46, 'topology_ieee39bus.png',
         '30–39 (Bus 39 = New England equiv.)',
         'Max H_B: Bus 1 (374.5 s)\nMin H_B: Bus 20 (26.0 s)\nSys strength: 2860 s\nAvg H_B: 73.3 s\nH_B ratio (max/min): 14.4×'),
        ('IEEE 68-Bus', 68, 16, 85, 'topology_ieee68bus.png',
         '1,9,22,31,46,54–62,65,68',
         'Max H_B: Bus 59 (500 s)\nMin H_B: Bus 56 (24.3 s)\nSys strength: 7201 s\nAvg H_B: 105.9 s\nH_B ratio (max/min): 20.6×'),
        ('IEEE 118-Bus', 118, 54, 186, 'topology_ieee118bus.png',
         '54 generators distributed',
         'Max H_B: Bus 4 (7.3 s)\nMin H_B: Bus 10 (4.2 s)\nSys strength: 652.5 s\nAvg H_B: 5.5 s\nH_B ratio (max/min): 1.7×'),
    ]
    for args in topo_data:
        slide_topology(prs, *args)
        print(f'  Slide: Topology {args[0]}')

    # Result slides
    result_pairs = [
        ('IEEE 14-Bus',  'inertia_bar_ieee14bus.png',  'freq_response_ieee14bus.png'),
        ('IEEE 39-Bus',  'inertia_bar_ieee39bus.png',  'freq_response_ieee39bus.png'),
        ('IEEE 68-Bus',  'inertia_bar_ieee68bus.png',  'freq_response_ieee68bus.png'),
        ('IEEE 118-Bus', 'inertia_bar_ieee118bus.png', 'freq_response_ieee118bus.png'),
    ]
    for name, bar, freq in result_pairs:
        slide_results_bar(prs, name, bar, freq)
        print(f'  Slide: Results {name}')

    slide_wc_heatmap(prs);     print('  Slide: W_c Heatmap')
    slide_res_comparison(prs); print('  Slide: RES Comparison')
    slide_analysis(prs);       print('  Slide: Analysis')
    slide_future(prs);         print('  Slide: Future Work')
    slide_conclusion(prs);     print('  Slide: Conclusion')

    prs.save(PPTX_OUT)
    print(f'\nPowerPoint saved: {PPTX_OUT}')
    print(f'Total slides: {len(prs.slides)}')


if __name__ == '__main__':
    main()
